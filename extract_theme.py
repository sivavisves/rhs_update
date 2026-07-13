from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE

def get_color(color_obj):
    if color_obj.type == MSO_COLOR_TYPE.RGB:
        return str(color_obj.rgb)
    elif color_obj.type == MSO_COLOR_TYPE.THEME:
        return f"Theme Color: {color_obj.theme_color}"
    return str(color_obj.type)

def extract_colors_from_pptx(file_path):
    prs = Presentation(file_path)
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, 'fill') and shape.fill.type is not None:
                try:
                    print(f"Shape Fill Color: {get_color(shape.fill.fore_color)}")
                except:
                    pass
            
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        print(f"Text: '{run.text.strip()}'")
                        if run.font.name:
                            print(f"  Font: {run.font.name}")
                        try:
                            if run.font.color and run.font.color.type:
                                print(f"  Color: {get_color(run.font.color)}")
                        except:
                            pass
             
if __name__ == '__main__':
    extract_colors_from_pptx('RHS Update.pptx')
